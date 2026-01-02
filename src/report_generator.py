"""
PPTX Report Generator for Famefact Social Media Reporting
Generates monthly reports from cached Facebook and Instagram data.
"""

import os
import io
import tempfile
import logging
from datetime import datetime, date
from typing import Optional, List, Dict, Any, Tuple
from dataclasses import dataclass
from pathlib import Path

import requests
from bs4 import BeautifulSoup
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

from .db import get_connection

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Famefact Brand Colors
FAMEFACT_DARK = RgbColor(0x1A, 0x1A, 0x2E)  # Dark background
FAMEFACT_CYAN = RgbColor(0x00, 0xD4, 0xFF)  # Primary accent
FAMEFACT_WHITE = RgbColor(0xFF, 0xFF, 0xFF)
FAMEFACT_GRAY = RgbColor(0x6B, 0x70, 0x80)
FAMEFACT_LIGHT_GRAY = RgbColor(0xF3, 0xF4, 0xF6)

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


@dataclass
class ReportConfig:
    """Configuration for report generation."""
    client_name: str
    report_month: str  # YYYY-MM format
    include_previous_months: bool = False
    fb_page_ids: Optional[List[str]] = None
    ig_account_ids: Optional[List[str]] = None
    output_dir: str = "/tmp/reports"


@dataclass
class PostData:
    """Data for a single post."""
    post_id: str
    page_id: str
    created_time: datetime
    post_type: str
    permalink: Optional[str]
    message: Optional[str]
    reactions_total: int
    comments_total: int
    shares_total: Optional[int]
    reach: Optional[int]
    impressions: Optional[int]
    video_3s_views: Optional[int]
    interactions_total: int
    media_url: Optional[str] = None
    thumbnail_url: Optional[str] = None


class ThumbnailHandler:
    """Handles fetching and caching of post thumbnails."""
    
    def __init__(self, cache_dir: Optional[str] = None):
        self.cache_dir = cache_dir or tempfile.mkdtemp(prefix="report_thumbnails_")
        Path(self.cache_dir).mkdir(parents=True, exist_ok=True)
        self.stats = {"graph_api": 0, "open_graph": 0, "missing": 0}
    
    def get_thumbnail(self, post: PostData) -> Optional[str]:
        """Get thumbnail path for a post."""
        # Try media_url first
        if post.media_url:
            path = self._download_image(post.media_url, post.post_id)
            if path:
                self.stats["graph_api"] += 1
                return path
        
        # Try thumbnail_url
        if post.thumbnail_url:
            path = self._download_image(post.thumbnail_url, post.post_id)
            if path:
                self.stats["graph_api"] += 1
                return path
        
        # Fallback to OpenGraph
        if post.permalink:
            og_url = self._get_og_image(post.permalink)
            if og_url:
                path = self._download_image(og_url, post.post_id)
                if path:
                    self.stats["open_graph"] += 1
                    return path
        
        self.stats["missing"] += 1
        return None
    
    def _download_image(self, url: str, post_id: str) -> Optional[str]:
        """Download image and return local path."""
        try:
            cache_path = os.path.join(self.cache_dir, f"{post_id}.jpg")
            if os.path.exists(cache_path):
                return cache_path
            
            response = requests.get(url, timeout=10, headers={
                "User-Agent": "Mozilla/5.0 (compatible; SocialDash/1.0)"
            })
            response.raise_for_status()
            
            # Convert to JPEG for consistency
            img = Image.open(io.BytesIO(response.content))
            if img.mode in ('RGBA', 'P'):
                img = img.convert('RGB')
            img.thumbnail((800, 800), Image.Resampling.LANCZOS)
            img.save(cache_path, "JPEG", quality=85)
            
            return cache_path
        except Exception as e:
            logger.warning(f"Failed to download image for {post_id}: {e}")
            return None
    
    def _get_og_image(self, permalink: str) -> Optional[str]:
        """Extract og:image from permalink."""
        try:
            response = requests.get(permalink, timeout=10, headers={
                "User-Agent": "Mozilla/5.0 (compatible; SocialDash/1.0)"
            })
            response.raise_for_status()
            
            soup = BeautifulSoup(response.text, 'html.parser')
            og_image = soup.find('meta', property='og:image')
            if og_image and og_image.get('content'):
                return og_image['content']
        except Exception as e:
            logger.warning(f"Failed to get OG image from {permalink}: {e}")
        return None


class ReportGenerator:
    """Generates PPTX reports from cached social media data."""
    
    def __init__(self, config: ReportConfig):
        self.config = config
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.thumbnail_handler = ThumbnailHandler()
        
        # Parse month
        year, month = map(int, config.report_month.split('-'))
        self.report_date = date(year, month, 1)
        self.month_name = self.report_date.strftime("%B %Y")
        self.month_name_de = self._get_german_month(month, year)
    
    def _get_german_month(self, month: int, year: int) -> str:
        """Get German month name."""
        months_de = [
            "Januar", "Februar", "März", "April", "Mai", "Juni",
            "Juli", "August", "September", "Oktober", "November", "Dezember"
        ]
        return f"{months_de[month - 1]} {year}"
    
    def generate(self) -> str:
        """Generate the complete report and return file path."""
        logger.info(f"Generating report for {self.config.client_name} - {self.month_name_de}")
        
        # Create slides
        self._add_cover_slide()
        
        # Facebook section
        fb_data = self._get_facebook_data()
        if fb_data:
            self._add_separator_slide("Facebook")
            self._add_facebook_kpi_slide(fb_data)
            self._add_facebook_top_posts_slide(fb_data)
            self._add_facebook_videos_slide(fb_data)
        
        # Instagram section
        ig_data = self._get_instagram_data()
        if ig_data:
            self._add_separator_slide("Instagram")
            self._add_instagram_kpi_slide(ig_data)
            self._add_instagram_top_posts_slide(ig_data)
            self._add_instagram_reels_slide(ig_data)
        
        # Fazit
        if fb_data or ig_data:
            self._add_fazit_slide(fb_data, ig_data)
        
        # Contact slide
        self._add_contact_slide()
        
        # Save
        output_path = self._save_report()
        
        # Log stats
        logger.info(f"Thumbnail stats: {self.thumbnail_handler.stats}")
        
        return output_path
    
    def _add_cover_slide(self):
        """Add cover slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Background
        self._set_slide_background(slide, FAMEFACT_DARK)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Social Media Reporting"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = FAMEFACT_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle (Month)
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(11.333), Inches(1)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = self.month_name_de
        p.font.size = Pt(32)
        p.font.color.rgb = FAMEFACT_CYAN
        p.alignment = PP_ALIGN.CENTER
        
        # Client name
        if self.config.client_name:
            client_box = slide.shapes.add_textbox(
                Inches(1), Inches(5), Inches(11.333), Inches(0.5)
            )
            tf = client_box.text_frame
            p = tf.paragraphs[0]
            p.text = self.config.client_name
            p.font.size = Pt(20)
            p.font.color.rgb = FAMEFACT_GRAY
            p.alignment = PP_ALIGN.CENTER
    
    def _add_separator_slide(self, title: str):
        """Add a section separator slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, FAMEFACT_DARK)
        
        # Platform icon/title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3), Inches(11.333), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = FAMEFACT_CYAN
        p.alignment = PP_ALIGN.CENTER
    
    def _add_facebook_kpi_slide(self, data: Dict[str, Any]):
        """Add Facebook KPI table slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, FAMEFACT_DARK)
        
        # Title
        self._add_slide_title(slide, "Facebook Kennzahlen")
        
        # KPI Table
        stats = data.get("stats", {})
        kpis = [
            ("Post-Reichweite", self._format_number(stats.get("total_reach", 0))),
            ("Ø Reichweite pro Post", self._format_number(stats.get("avg_reach_per_post", 0))),
            ("Interaktionen", self._format_number(stats.get("total_interactions", 0))),
            ("Reactions", self._format_number(stats.get("total_reactions", 0))),
            ("Kommentare", self._format_number(stats.get("total_comments", 0))),
            ("Video Views (3-Sek)", self._format_number(stats.get("total_video_views", 0))),
            ("Anzahl Postings", str(stats.get("total_posts", 0))),
        ]
        
        # Calculate interaction rate
        if stats.get("total_reach", 0) > 0:
            rate = (stats.get("total_interactions", 0) / stats.get("total_reach", 0)) * 100
            kpis.append(("Interaktionsrate", f"{rate:.2f}%"))
        
        # Add shares separately
        shares = stats.get("total_shares", 0)
        kpis.append(("Shares (Limited)", self._format_number(shares)))
        
        self._add_kpi_table(slide, kpis, self.month_name_de)
        
        # Footnote
        self._add_footnote(slide, [
            "Interaktionsrate = Interaktionen / Post-Reichweite × 100",
            "Interaktionen = Reactions + Comments (Shares separat, da eingeschränkt)"
        ])
    
    def _add_facebook_top_posts_slide(self, data: Dict[str, Any]):
        """Add Facebook top posts slide with images."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, FAMEFACT_DARK)
        self._add_slide_title(slide, "Postings nach Interaktion – Bilder")
        
        posts = data.get("top_image_posts", [])[:6]
        self._add_post_grid(slide, posts, "fb")
        
        # Footnote
        self._add_footnote(slide, [
            "In die Interaktionen fallen Reactions und Kommentare. Shares separat (limited)."
        ])
    
    def _add_facebook_videos_slide(self, data: Dict[str, Any]):
        """Add Facebook videos slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, FAMEFACT_DARK)
        self._add_slide_title(slide, "Videos nach 3-sekündige Video Views")
        
        videos = data.get("top_videos", [])[:6]
        if videos:
            self._add_post_grid(slide, videos, "fb", show_views=True)
        else:
            self._add_no_data_message(slide, "Keine Video-Daten für diesen Monat verfügbar.")
    
    def _add_instagram_kpi_slide(self, data: Dict[str, Any]):
        """Add Instagram KPI table slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, FAMEFACT_DARK)
        self._add_slide_title(slide, "Instagram Kennzahlen")
        
        stats = data.get("stats", {})
        kpis = [
            ("Post-Reichweite", self._format_number(stats.get("total_reach", 0))),
            ("Ø Reichweite pro Post", self._format_number(stats.get("avg_reach_per_post", 0))),
            ("Interaktionen", self._format_number(stats.get("total_interactions", 0))),
            ("Likes", self._format_number(stats.get("total_likes", 0))),
            ("Kommentare", self._format_number(stats.get("total_comments", 0))),
            ("Saves", self._format_number(stats.get("total_saves", 0))),
            ("Video/Reels Plays", self._format_number(stats.get("total_plays", 0))),
            ("Anzahl Postings", str(stats.get("total_posts", 0))),
        ]
        
        # Calculate engagement rate
        if stats.get("total_reach", 0) > 0:
            rate = (stats.get("total_interactions", 0) / stats.get("total_reach", 0)) * 100
            kpis.append(("Engagement Rate", f"{rate:.2f}%"))
        
        self._add_kpi_table(slide, kpis, self.month_name_de)
        
        # Footnote
        self._add_footnote(slide, [
            "Interaktionen = Likes + Kommentare + Saves",
            "Engagement Rate = Interaktionen / Reichweite × 100"
        ])
    
    def _add_instagram_top_posts_slide(self, data: Dict[str, Any]):
        """Add Instagram top posts slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, FAMEFACT_DARK)
        self._add_slide_title(slide, "Postings nach Interaktion – Bilder")
        
        posts = data.get("top_image_posts", [])[:6]
        self._add_post_grid(slide, posts, "ig")
    
    def _add_instagram_reels_slide(self, data: Dict[str, Any]):
        """Add Instagram reels slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, FAMEFACT_DARK)
        self._add_slide_title(slide, "Reels: Aufrufe/Plays")
        
        reels = data.get("top_reels", [])[:6]
        if reels:
            self._add_post_grid(slide, reels, "ig", show_plays=True)
        else:
            self._add_no_data_message(slide, "Keine Reels-Daten für diesen Monat verfügbar.")
    
    def _add_fazit_slide(self, fb_data: Optional[Dict], ig_data: Optional[Dict]):
        """Add conclusion slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, FAMEFACT_DARK)
        self._add_slide_title(slide, "Fazit")
        
        fazit_text = []
        
        # Facebook summary
        if fb_data:
            stats = fb_data.get("stats", {})
            fb_summary = f"Facebook: Im {self.month_name_de} wurden {stats.get('total_posts', 0)} Beiträge veröffentlicht. "
            fb_summary += f"Die Gesamtreichweite betrug {self._format_number(stats.get('total_reach', 0))} "
            fb_summary += f"mit {self._format_number(stats.get('total_interactions', 0))} Interaktionen. "
            
            top_post = fb_data.get("top_image_posts", [{}])[0] if fb_data.get("top_image_posts") else None
            if top_post:
                fb_summary += f"Der Top-Post erzielte {self._format_number(top_post.get('interactions_total', 0))} Interaktionen."
            
            fazit_text.append(fb_summary)
        
        # Instagram summary
        if ig_data:
            stats = ig_data.get("stats", {})
            ig_summary = f"Instagram: Im {self.month_name_de} wurden {stats.get('total_posts', 0)} Beiträge veröffentlicht. "
            ig_summary += f"Die Gesamtreichweite betrug {self._format_number(stats.get('total_reach', 0))} "
            ig_summary += f"mit {self._format_number(stats.get('total_interactions', 0))} Interaktionen "
            ig_summary += f"und {self._format_number(stats.get('total_saves', 0))} Saves. "
            
            top_reel = ig_data.get("top_reels", [{}])[0] if ig_data.get("top_reels") else None
            if top_reel:
                ig_summary += f"Das Top-Reel erzielte {self._format_number(top_reel.get('plays', 0))} Plays."
            
            fazit_text.append(ig_summary)
        
        # Add text to slide
        text_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.5), Inches(11.833), Inches(5)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for i, text in enumerate(fazit_text):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(16)
            p.font.color.rgb = FAMEFACT_WHITE
            p.space_after = Pt(24)
    
    def _add_contact_slide(self):
        """Add contact/signature slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, FAMEFACT_DARK)
        
        # Thank you text
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Vielen Dank"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = FAMEFACT_CYAN
        p.alignment = PP_ALIGN.CENTER
        
        # Contact info
        contact_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(11.333), Inches(1)
        )
        tf = contact_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Report erstellt für {self.config.client_name}"
        p.font.size = Pt(18)
        p.font.color.rgb = FAMEFACT_GRAY
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = f"{self.month_name_de}"
        p.font.size = Pt(14)
        p.font.color.rgb = FAMEFACT_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    def _set_slide_background(self, slide, color: RgbColor):
        """Set slide background color."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
    
    def _add_slide_title(self, slide, title: str):
        """Add title to slide."""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = FAMEFACT_WHITE
    
    def _add_kpi_table(self, slide, kpis: List[Tuple[str, str]], month_label: str):
        """Add KPI table to slide."""
        rows = len(kpis) + 1  # +1 for header
        cols = 2
        
        table = slide.shapes.add_table(
            rows, cols,
            Inches(1), Inches(1.3),
            Inches(8), Inches(0.4 * rows)
        ).table
        
        # Header
        table.cell(0, 0).text = "KPI"
        table.cell(0, 1).text = month_label
        
        # Style header
        for col in range(cols):
            cell = table.cell(0, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = FAMEFACT_CYAN
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = FAMEFACT_DARK
        
        # Data rows
        for i, (kpi_name, kpi_value) in enumerate(kpis, start=1):
            table.cell(i, 0).text = kpi_name
            table.cell(i, 1).text = kpi_value
            
            for col in range(cols):
                cell = table.cell(i, col)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(0x2D, 0x2D, 0x44) if i % 2 == 0 else RgbColor(0x1A, 0x1A, 0x2E)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(11)
                    paragraph.font.color.rgb = FAMEFACT_WHITE
    
    def _add_post_grid(self, slide, posts: List[Dict], platform: str, 
                       show_views: bool = False, show_plays: bool = False):
        """Add post grid with thumbnails."""
        cols = 3
        rows = 2
        
        start_x = Inches(0.5)
        start_y = Inches(1.3)
        card_width = Inches(4)
        card_height = Inches(2.8)
        gap = Inches(0.2)
        
        for i, post in enumerate(posts[:6]):
            row = i // cols
            col = i % cols
            
            x = start_x + col * (card_width + gap)
            y = start_y + row * (card_height + gap)
            
            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RgbColor(0x2D, 0x2D, 0x44)
            card.line.fill.background()
            
            # Try to add thumbnail
            post_data = self._dict_to_post_data(post, platform)
            thumbnail_path = self.thumbnail_handler.get_thumbnail(post_data)
            
            if thumbnail_path:
                try:
                    slide.shapes.add_picture(
                        thumbnail_path,
                        x + Inches(0.1), y + Inches(0.1),
                        width=card_width - Inches(0.2),
                        height=Inches(1.8)
                    )
                except Exception as e:
                    logger.warning(f"Failed to add image: {e}")
                    self._add_placeholder_image(slide, x + Inches(0.1), y + Inches(0.1), 
                                                card_width - Inches(0.2), Inches(1.8))
            else:
                self._add_placeholder_image(slide, x + Inches(0.1), y + Inches(0.1), 
                                            card_width - Inches(0.2), Inches(1.8))
            
            # Date
            date_str = post.get('post_created_time', post.get('created_time', ''))
            if isinstance(date_str, datetime):
                date_str = date_str.strftime("%d.%m.%Y")
            elif isinstance(date_str, str) and date_str:
                try:
                    date_str = datetime.fromisoformat(date_str.replace('Z', '+00:00')).strftime("%d.%m.%Y")
                except:
                    pass
            
            date_box = slide.shapes.add_textbox(
                x + Inches(0.1), y + Inches(2),
                card_width - Inches(0.2), Inches(0.3)
            )
            tf = date_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(date_str)
            p.font.size = Pt(9)
            p.font.color.rgb = FAMEFACT_GRAY
            
            # Metrics
            metrics_box = slide.shapes.add_textbox(
                x + Inches(0.1), y + Inches(2.3),
                card_width - Inches(0.2), Inches(0.4)
            )
            tf = metrics_box.text_frame
            p = tf.paragraphs[0]
            
            if show_views:
                views = post.get('video_3s_views', 0) or 0
                p.text = f"👁 {self._format_number(views)} Views"
            elif show_plays:
                plays = post.get('plays', 0) or 0
                p.text = f"▶ {self._format_number(plays)} Plays"
            else:
                interactions = post.get('interactions_total', 0) or 0
                p.text = f"❤ {self._format_number(interactions)} Interaktionen"
            
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = FAMEFACT_CYAN
    
    def _add_placeholder_image(self, slide, x, y, width, height):
        """Add placeholder for missing image."""
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, height
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RgbColor(0x3D, 0x3D, 0x54)
        placeholder.line.fill.background()
        
        # Add text
        text_box = slide.shapes.add_textbox(x, y + height / 2 - Inches(0.2), width, Inches(0.4))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Preview nicht verfügbar"
        p.font.size = Pt(10)
        p.font.color.rgb = FAMEFACT_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    def _add_footnote(self, slide, notes: List[str]):
        """Add footnote to slide."""
        footnote_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.8),
            Inches(12.333), Inches(0.5)
        )
        tf = footnote_box.text_frame
        
        for i, note in enumerate(notes):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = note
            p.font.size = Pt(8)
            p.font.color.rgb = FAMEFACT_GRAY
    
    def _add_no_data_message(self, slide, message: str):
        """Add no data message to slide."""
        msg_box = slide.shapes.add_textbox(
            Inches(1), Inches(3), Inches(11.333), Inches(1)
        )
        tf = msg_box.text_frame
        p = tf.paragraphs[0]
        p.text = message
        p.font.size = Pt(18)
        p.font.color.rgb = FAMEFACT_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    def _format_number(self, num) -> str:
        """Format number for display."""
        if num is None:
            return "0"
        num = int(num)
        if num >= 1_000_000:
            return f"{num / 1_000_000:.1f}M"
        if num >= 1_000:
            return f"{num / 1_000:.1f}K"
        return f"{num:,}".replace(",", ".")
    
    def _dict_to_post_data(self, post: Dict, platform: str) -> PostData:
        """Convert dict to PostData."""
        if platform == "fb":
            return PostData(
                post_id=post.get('post_id', ''),
                page_id=post.get('page_id', ''),
                created_time=post.get('post_created_time', datetime.now()),
                post_type=post.get('post_type', ''),
                permalink=post.get('permalink'),
                message=post.get('message'),
                reactions_total=post.get('reactions_total', 0) or 0,
                comments_total=post.get('comments_total', 0) or 0,
                shares_total=post.get('shares_total'),
                reach=post.get('reach'),
                impressions=post.get('impressions'),
                video_3s_views=post.get('video_3s_views'),
                interactions_total=post.get('interactions_total', 0) or 0,
                media_url=post.get('media_url'),
                thumbnail_url=post.get('thumbnail_url'),
            )
        else:  # ig
            return PostData(
                post_id=post.get('media_id', ''),
                page_id=post.get('account_id', ''),
                created_time=post.get('post_created_time', datetime.now()),
                post_type=post.get('media_type', ''),
                permalink=post.get('permalink'),
                message=post.get('caption'),
                reactions_total=post.get('likes', 0) or 0,
                comments_total=post.get('comments', 0) or 0,
                shares_total=post.get('shares'),
                reach=post.get('reach'),
                impressions=post.get('impressions'),
                video_3s_views=post.get('plays'),
                interactions_total=post.get('interactions_total', 0) or 0,
                media_url=post.get('media_url'),
                thumbnail_url=post.get('thumbnail_url'),
            )
    
    def _get_facebook_data(self) -> Optional[Dict[str, Any]]:
        """Fetch Facebook data from database."""
        conn = get_connection()
        if not conn:
            logger.warning("No database connection for Facebook data")
            return None
        
        try:
            cur = conn.cursor()
            month_str = self.report_date.strftime("%Y-%m-%d")
            
            # Get monthly stats
            cur.execute("""
                SELECT * FROM view_fb_monthly_page_stats 
                WHERE month = %s
            """, (month_str,))
            
            stats_row = cur.fetchone()
            if not stats_row:
                logger.info("No Facebook stats for this month")
                return None
            
            columns = [desc[0] for desc in cur.description]
            stats = dict(zip(columns, stats_row))
            
            # Get top image posts
            cur.execute("""
                SELECT * FROM view_fb_monthly_post_metrics 
                WHERE month = %s 
                AND post_type IN ('photo', 'image', 'link')
                ORDER BY interactions_total DESC
                LIMIT 9
            """, (month_str,))
            
            top_image_posts = []
            columns = [desc[0] for desc in cur.description]
            for row in cur.fetchall():
                top_image_posts.append(dict(zip(columns, row)))
            
            # Get top videos
            cur.execute("""
                SELECT * FROM view_fb_monthly_post_metrics 
                WHERE month = %s 
                AND post_type IN ('video', 'reel')
                AND video_3s_views IS NOT NULL
                ORDER BY video_3s_views DESC NULLS LAST
                LIMIT 6
            """, (month_str,))
            
            top_videos = []
            columns = [desc[0] for desc in cur.description]
            for row in cur.fetchall():
                top_videos.append(dict(zip(columns, row)))
            
            cur.close()
            conn.close()
            
            return {
                "stats": stats,
                "top_image_posts": top_image_posts,
                "top_videos": top_videos
            }
            
        except Exception as e:
            logger.error(f"Error fetching Facebook data: {e}")
            return None
    
    def _get_instagram_data(self) -> Optional[Dict[str, Any]]:
        """Fetch Instagram data from database."""
        conn = get_connection()
        if not conn:
            logger.warning("No database connection for Instagram data")
            return None
        
        try:
            cur = conn.cursor()
            month_str = self.report_date.strftime("%Y-%m-%d")
            
            # Get monthly stats
            cur.execute("""
                SELECT * FROM view_ig_monthly_account_stats 
                WHERE month = %s
            """, (month_str,))
            
            stats_row = cur.fetchone()
            if not stats_row:
                logger.info("No Instagram stats for this month")
                return None
            
            columns = [desc[0] for desc in cur.description]
            stats = dict(zip(columns, stats_row))
            
            # Get top image posts
            cur.execute("""
                SELECT * FROM view_ig_monthly_post_metrics 
                WHERE month = %s 
                AND media_type IN ('IMAGE', 'CAROUSEL_ALBUM')
                ORDER BY interactions_total DESC
                LIMIT 9
            """, (month_str,))
            
            top_image_posts = []
            columns = [desc[0] for desc in cur.description]
            for row in cur.fetchall():
                top_image_posts.append(dict(zip(columns, row)))
            
            # Get top reels
            cur.execute("""
                SELECT * FROM view_ig_monthly_post_metrics 
                WHERE month = %s 
                AND media_type = 'REEL'
                ORDER BY plays DESC NULLS LAST
                LIMIT 6
            """, (month_str,))
            
            top_reels = []
            columns = [desc[0] for desc in cur.description]
            for row in cur.fetchall():
                top_reels.append(dict(zip(columns, row)))
            
            cur.close()
            conn.close()
            
            return {
                "stats": stats,
                "top_image_posts": top_image_posts,
                "top_reels": top_reels
            }
            
        except Exception as e:
            logger.error(f"Error fetching Instagram data: {e}")
            return None
    
    def _save_report(self) -> str:
        """Save the report and return file path."""
        Path(self.config.output_dir).mkdir(parents=True, exist_ok=True)
        
        filename = f"report_{self.config.client_name.lower().replace(' ', '_')}_{self.config.report_month}.pptx"
        filepath = os.path.join(self.config.output_dir, filename)
        
        self.prs.save(filepath)
        logger.info(f"Report saved to: {filepath}")
        
        return filepath


def generate_report(
    client_name: str,
    report_month: str,
    output_dir: str = "/tmp/reports",
    include_previous_months: bool = False
) -> str:
    """
    Generate a PPTX report for the specified month.
    
    Args:
        client_name: Name of the client/brand
        report_month: Month in YYYY-MM format
        output_dir: Directory to save the report
        include_previous_months: Include comparison with previous months
    
    Returns:
        Path to the generated PPTX file
    """
    config = ReportConfig(
        client_name=client_name,
        report_month=report_month,
        output_dir=output_dir,
        include_previous_months=include_previous_months
    )
    
    generator = ReportGenerator(config)
    return generator.generate()


if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="Generate Social Media Report")
    parser.add_argument("--client", required=True, help="Client name")
    parser.add_argument("--month", required=True, help="Report month (YYYY-MM)")
    parser.add_argument("--output", default="/tmp/reports", help="Output directory")
    parser.add_argument("--include-previous", action="store_true", help="Include previous months")
    
    args = parser.parse_args()
    
    report_path = generate_report(
        client_name=args.client,
        report_month=args.month,
        output_dir=args.output,
        include_previous_months=args.include_previous
    )
    
    print(f"Report generated: {report_path}")
